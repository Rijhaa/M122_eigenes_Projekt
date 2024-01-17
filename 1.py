from docx import Document
from pptx import Presentation
from pptx.util import Pt

def word_to_powerpoint(word_file, ppt_file, master_pptx):
    # Word-Dokument öffnen
    doc = Document(word_file)

    # PowerPoint-Präsentation erstellen, basierend auf dem Master-Design
    presentation = Presentation(master_pptx)

    # Benutzerdefinierte Funktion zur Erstellung einer Folie basierend auf dem Master-Design
    def create_slide(title, content):
        # Schriftgröße und -art einstellen
        font_size = Pt(12)
        font_name = 'Arial'

        # Neue Folie erstellen (basierend auf Master-Folie)
        slide_layout = presentation.slide_layouts[0]  # 0 entspricht der Master-Folie
        slide = presentation.slides.add_slide(slide_layout)

        # Iteriere über alle Formen in der Folie und ändere die Schriftgröße und -art
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = font_size
                        run.font.name = font_name

        # Setze den Titel
        title_shape = slide.shapes.title
        title_shape.text_frame.text = title

        # Setze den Inhalt
        content_shape = slide.placeholders[1]
        content_shape.text_frame.text = content

    title = None  # Variable für den aktuellen Titel initialisieren

    # Iteriere durch alle Absätze im Word-Dokument
    for paragraph in doc.paragraphs:
        # Überprüfe, ob der Absatz kursiv formatiert ist
        if paragraph.runs and paragraph.runs[0].italic:
            title = paragraph.text.strip()  # Titel aktualisieren

        # Nur nicht-leeren Text in die Folie einfügen
        elif paragraph.text.strip():
            # Wenn es keinen expliziten Titel gibt, verwende einen Standardtitel
            if title is None:
                title = "Titel hier einfügen"

            create_slide(title, paragraph.text)

            # Nach dem Erstellen der Folie den Titel zurücksetzen
            title = None

    # PowerPoint-Präsentation speichern
    presentation.save(ppt_file)

# Beispielaufruf der Funktion
if __name__ == "__main__":
    # Ersetzen Sie die Pfade durch die tatsächlichen Dateipfade Ihrer Word-, PowerPoint- und Master-Design-Dateien
    word_to_powerpoint("Text_Datenbank.docx", "Präsentation_Datenbank.pptx", "MasterDesign.pptx")
